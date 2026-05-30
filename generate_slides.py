"""
Sinh slide thuyết trình PowerPoint cho đồ án EZWare.
Mục tiêu: 15 slide, thuyết trình trong 15 phút (1 phút/slide).

Color palette: Midnight Executive
  - Primary navy: #112D4E
  - Mid blue:     #3F72AF
  - Light blue:   #DBE2EF
  - Off-white:    #F9F7F7
  - Accent:       #DC5F00 (cam dùng cho điểm nhấn / chức năng quan trọng)

Re-use 6 hình từ thư mục report_images/ (đã sinh từ generate_diagrams.py).

Chạy:
    python generate_slides.py
Output: NhomXX_Slide_ThuyetTrinh.pptx
"""
from __future__ import annotations
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


OUTPUT = 'NhomXX_Slide_ThuyetTrinh.pptx'
IMG_DIR = 'report_images'
DEPLOY_URL = 'https://ezware-backend.onrender.com'

# Color palette — Midnight Executive
NAVY = RGBColor(0x11, 0x2D, 0x4E)
MID = RGBColor(0x3F, 0x72, 0xAF)
LIGHT = RGBColor(0xDB, 0xE2, 0xEF)
OFFWHITE = RGBColor(0xF9, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xDC, 0x5F, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)


# ============================================================
# Helpers
# ============================================================
def hex_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill, line=None, corner=0.1):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = corner
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=DARK_GRAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_circle(slide, cx, cy, d, fill, line=None):
    """Vòng tròn (toạ độ theo center)."""
    r = d / 2
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - r), Inches(cy - r),
                                 Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_number_footer(slide, idx, total, dark_bg=False):
    color = OFFWHITE if dark_bg else GRAY
    add_text(slide, 12.5, 7.15, 0.6, 0.3,
             f'{idx} / {total}', size=10, color=color, align=PP_ALIGN.RIGHT)


def add_title_bar(slide, title, subtitle=None):
    """Title chuẩn — không dùng accent line dưới (AI slop)."""
    add_text(slide, 0.5, 0.35, 12.3, 0.7,
             title, size=30, bold=True, color=NAVY, font='Calibri')
    if subtitle:
        add_text(slide, 0.5, 1.0, 12.3, 0.4,
                 subtitle, size=14, italic=True, color=MID, font='Calibri')


# ============================================================
# Build slide
# ============================================================
def build_cover(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(s, NAVY)

    # Vòng tròn decorative top-right
    add_circle(s, 12.5, 1.0, 1.6, MID)
    add_circle(s, 11.7, 1.5, 0.8, ACCENT)

    # Brand
    add_text(s, 0.6, 0.6, 6, 0.4,
             'ĐẠI HỌC CÔNG NGHỆ THÔNG TIN — ĐHQG TP.HCM',
             size=12, color=LIGHT)
    add_text(s, 0.6, 1.0, 6, 0.35,
             'ĐỒ ÁN MÔN HỌC IE221 — LẬP TRÌNH PYTHON',
             size=11, italic=True, color=LIGHT)

    # Main title — center
    add_text(s, 0.6, 2.4, 12.1, 0.9,
             'HỆ THỐNG QUẢN LÝ KHO EZWARE',
             size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
             font='Calibri')
    add_text(s, 0.6, 3.3, 12.1, 0.6,
             'Xây dựng Backend API với Django REST Framework',
             size=22, color=LIGHT, font='Calibri')

    # Accent dot
    add_circle(s, 1.0, 4.4, 0.15, ACCENT)
    add_text(s, 1.25, 4.27, 11.5, 0.4,
             'Triển khai công khai tại  ezware-backend.onrender.com/swagger/',
             size=14, italic=True, color=ACCENT, font='Consolas')

    # Group + members box
    box = add_round(s, 0.6, 5.3, 7.5, 1.7, MID, corner=0.05)
    add_text(s, 0.85, 5.4, 7, 0.35,
             'NHÓM XX',
             size=18, bold=True, color=WHITE, font='Calibri')
    add_text(s, 0.85, 5.75, 7, 1.2, [
        'Sinh viên 1:  Nguyễn Văn A  —  25XXXXX1',
        'Sinh viên 2:  Nguyễn Văn B  —  25XXXXX2',
        'GVHD:  …',
    ], size=14, color=WHITE, font='Calibri')

    add_text(s, 8.5, 6.3, 4, 0.4,
             f'TP.HCM — 05/2026',
             size=14, italic=True, color=LIGHT,
             align=PP_ALIGN.RIGHT, font='Calibri')


def build_agenda(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, 'Nội dung thuyết trình', 'Khoảng 15 phút')

    items = [
        ('1', 'Giới thiệu đề tài', 'Vấn đề và giải pháp'),
        ('2', 'Mục tiêu & Phạm vi', '4 yêu cầu cốt lõi'),
        ('3', 'Công nghệ sử dụng', 'Python · Django · DRF'),
        ('4', 'Kiến trúc MVT', 'Tổ chức 6 app Django'),
        ('5', 'Cơ sở dữ liệu', '6 bảng + ràng buộc'),
        ('6', 'Chức năng hệ thống', '4 nhóm chính'),
        ('7', '14 URL endpoint REST', '23 thao tác HTTP'),
        ('8', 'Demo nghiệp vụ then chốt', 'Duyệt phiếu — atomic'),
        ('9', 'Kiểm thử Postman + Swagger', '30+ test case'),
        ('10', 'Deploy thực tế trên Render', 'Public Internet'),
        ('11', 'Thách thức kỹ thuật đã giải quyết', 'OOM, pkg_resources'),
        ('12', 'Kết luận & Hướng phát triển', ''),
    ]
    # Layout 2 cột
    col_w = 6.0
    row_h = 0.46
    start_y = 1.6
    for i, (num, title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.5 + col * (col_w + 0.3)
        y = start_y + row * row_h

        add_circle(s, x + 0.25, y + 0.23, 0.35, MID)
        add_text(s, x + 0.08, y + 0.05, 0.4, 0.36,
                 num, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font='Calibri')
        add_text(s, x + 0.55, y + 0.02, col_w - 0.65, 0.22,
                 title, size=13, bold=True, color=NAVY, font='Calibri')
        if desc:
            add_text(s, x + 0.55, y + 0.24, col_w - 0.65, 0.22,
                     desc, size=10, color=GRAY, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_intro(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '1. Giới thiệu đề tài',
                  'Bối cảnh — Vấn đề — Giải pháp')

    # 3 cột: Bối cảnh / Vấn đề / Giải pháp
    headers = [
        ('BỐI CẢNH', MID,
         'Doanh nghiệp nhỏ và vừa cần một hệ thống quản lý kho cơ bản nhưng có '
         'business logic chuẩn (phiếu nhập, phiếu xuất, tồn kho cập nhật chính '
         'xác — kể cả khi nhiều người duyệt cùng lúc).'),
        ('VẤN ĐỀ', ACCENT,
         'Các giải pháp ERP thương mại đắt và phức tạp. Tự code thường gặp lỗi '
         'race condition khi cập nhật tồn kho song song; lỗi phân quyền cho '
         'phép leo thang quyền (privilege escalation).'),
        ('GIẢI PHÁP CỦA NHÓM', NAVY,
         'Xây dựng backend REST API “EZWare” bằng Django + DRF, bảo vệ tồn kho '
         'bằng transaction.atomic + select_for_update, phân quyền theo role '
         'admin/staff, sinh tài liệu tự động bằng Swagger.'),
    ]
    col_w = 4.1
    col_x = [0.5, 4.7, 8.9]
    for (head, color, body), x in zip(headers, col_x):
        # Header bar
        add_round(s, x, 1.6, col_w, 0.55, color, corner=0.15)
        add_text(s, x, 1.65, col_w, 0.45,
                 head, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font='Calibri')
        # Body
        add_round(s, x, 2.25, col_w, 3.6, OFFWHITE, line=LIGHT, corner=0.03)
        add_text(s, x + 0.2, 2.4, col_w - 0.4, 3.3,
                 body, size=13, color=DARK_GRAY, font='Calibri')

    # Stats bar bottom
    stats = [
        ('14', 'URL endpoint REST'),
        ('23', 'HTTP operations'),
        ('6', 'Bảng CSDL'),
        ('30+', 'Postman test case'),
    ]
    stat_w = 2.9
    for i, (val, lab) in enumerate(stats):
        x = 0.6 + i * (stat_w + 0.2)
        add_round(s, x, 6.2, stat_w, 0.85, NAVY, corner=0.1)
        add_text(s, x, 6.27, stat_w, 0.45,
                 val, size=24, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER, font='Calibri')
        add_text(s, x, 6.72, stat_w, 0.3,
                 lab, size=10, color=LIGHT,
                 align=PP_ALIGN.CENTER, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_goals(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '2. Mục tiêu & Phạm vi đề tài',
                  '4 yêu cầu cốt lõi đã được đáp ứng đầy đủ')

    goals = [
        ('CSDL', '6 bảng quan hệ',
         'Phục vụ nghiệp vụ: user, sản phẩm, kho, phiếu (header + detail), tồn '
         'kho. Tiếng Việt có dấu, không dùng ngoại tệ.'),
        ('NGƯỜI DÙNG', 'Tối thiểu 2 vai trò',
         'admin (quản trị) và staff (nhân viên kho) — phân quyền cứng bằng '
         'permission class.'),
        ('REST API', 'GET / POST / PUT / DELETE',
         '14 URL endpoint với 23 thao tác HTTP, tài liệu tự sinh bằng Swagger UI, '
         'xác thực Token.'),
        ('KIỂM THỬ', 'Postman + script',
         'Collection 30+ test case kèm script auto-save token, regression test '
         'cho lỗ hổng leo thang quyền.'),
    ]

    box_w = 6.0
    box_h = 2.4
    positions = [(0.5, 1.6), (6.85, 1.6), (0.5, 4.2), (6.85, 4.2)]

    for (head, sub, body), (x, y) in zip(goals, positions):
        add_round(s, x, y, box_w, box_h, OFFWHITE, line=LIGHT, corner=0.02)
        # Icon circle
        add_circle(s, x + 0.6, y + 0.55, 0.8, NAVY)
        add_text(s, x + 0.2, y + 0.2, 0.8, 0.7,
                 head[0], size=22, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Calibri')
        add_text(s, x + 1.15, y + 0.15, box_w - 1.3, 0.4,
                 head, size=16, bold=True, color=NAVY, font='Calibri')
        add_text(s, x + 1.15, y + 0.55, box_w - 1.3, 0.4,
                 sub, size=13, italic=True, color=MID, font='Calibri')
        add_text(s, x + 0.3, y + 1.1, box_w - 0.5, box_h - 1.2,
                 body, size=13, color=DARK_GRAY, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_tech(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '3. Công nghệ sử dụng',
                  'Stack Python thuần — production-ready')

    tech_groups = [
        ('NGÔN NGỮ & FRAMEWORK', NAVY, [
            'Python 3.10+ (production 3.12.3)',
            'Django 5.0.6',
            'Django REST Framework 3.15.1',
        ]),
        ('THƯ VIỆN HỖ TRỢ', MID, [
            'drf-yasg 1.21.7 (Swagger UI tự sinh)',
            'pandas + openpyxl (import_excel)',
            'gunicorn 23.0.0 + whitenoise 6.7.0 (production)',
        ]),
        ('CƠ SỞ DỮ LIỆU & DEPLOY', ACCENT, [
            'SQLite (file đơn lẻ, không cần server)',
            'Render.com — Python 3.12 web service (free tier)',
            'GitHub — quản lý source code',
        ]),
    ]
    block_w = 12.3
    block_h = 1.5   # giảm từ 1.65 → 1.5 để chừa chỗ cho caption bottom
    for i, (head, color, items) in enumerate(tech_groups):
        y = 1.55 + i * (block_h + 0.12)
        # Header strip bên trái
        add_round(s, 0.5, y, 3.4, block_h, color, corner=0.05)
        add_text(s, 0.5, y + 0.1, 3.4, 0.5,
                 head, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Calibri')
        # Body bên phải
        add_round(s, 4.0, y, 8.8, block_h, OFFWHITE, line=LIGHT, corner=0.02)
        bullet_text = '\n'.join(f'•  {it}' for it in items)
        add_text(s, 4.2, y + 0.1, 8.4, block_h - 0.15,
                 bullet_text.split('\n'),
                 size=13, color=DARK_GRAY, font='Calibri')

    # Bottom note — đẩy xuống y=6.95 để không đè lên block cuối (end ở y=6.4)
    add_text(s, 0.5, 6.95, 12.3, 0.4,
             'Tách prod/dev requirements: pandas chỉ ở local cho import data — '
             'giúp build trên Render không OOM 512 MB RAM.',
             size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_architecture(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '4. Kiến trúc hệ thống — Mô hình MVT',
                  'Biến thể MVC của Django, áp dụng OOP triệt để')

    # Bên trái: sơ đồ MVT 3 tầng
    layers = [
        ('CLIENT', 'Postman / Browser / curl', LIGHT, NAVY),
        ('VIEW (Class-based)', 'ListCreateAPIView, RetrieveUpdateDestroyAPIView, APIView', MID, WHITE),
        ('SERIALIZER (DRF)', 'ModelSerializer ↔ JSON', ACCENT, WHITE),
        ('MODEL (ORM)', '6 lớp Django Model ↔ 6 bảng SQLite', NAVY, WHITE),
        ('DATABASE', 'SQLite', LIGHT, NAVY),
    ]
    layer_w = 6.0
    layer_h = 0.75
    start_y = 1.6
    for i, (name, desc, bg, fg) in enumerate(layers):
        y = start_y + i * (layer_h + 0.15)
        add_round(s, 0.5, y, layer_w, layer_h, bg, corner=0.06)
        add_text(s, 0.7, y + 0.05, layer_w - 0.3, 0.32,
                 name, size=14, bold=True, color=fg, font='Calibri')
        add_text(s, 0.7, y + 0.38, layer_w - 0.3, 0.3,
                 desc, size=11, italic=True, color=fg, font='Calibri')
        # Arrow down between layers
        if i < len(layers) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(3.3), Inches(y + layer_h + 0.0),
                                     Inches(0.4), Inches(0.15))
            arr.fill.solid(); arr.fill.fore_color.rgb = GRAY
            arr.line.fill.background()

    # Bên phải: 6 Django apps
    add_text(s, 7.0, 1.55, 5.8, 0.4,
             '6 Django Apps (tách theo domain)',
             size=15, bold=True, color=NAVY, font='Calibri')

    apps = [
        ('accounts', 'User, Auth, Token, Permission'),
        ('products', 'Danh mục sản phẩm'),
        ('warehouses', 'Danh mục kho'),
        ('inventory', 'Phiếu nhập/xuất, tồn kho'),
        ('reports', 'Báo cáo tồn thấp'),
        ('core', 'Constants + import_excel'),
    ]
    for i, (name, desc) in enumerate(apps):
        col = i % 2
        row = i // 2
        x = 7.0 + col * 3.0
        y = 2.05 + row * 1.1

        add_round(s, x, y, 2.85, 0.95, OFFWHITE, line=MID, corner=0.05)
        # Tag
        add_round(s, x + 0.1, y + 0.1, 1.3, 0.3, NAVY, corner=0.3)
        add_text(s, x + 0.1, y + 0.1, 1.3, 0.3,
                 name, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Consolas')
        add_text(s, x + 0.1, y + 0.45, 2.7, 0.5,
                 desc, size=10, color=DARK_GRAY, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_database(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '5. Cơ sở dữ liệu',
                  '6 bảng SQLite + 5 ràng buộc khóa ngoại — match yêu cầu 5-10 bảng')

    # Insert ER diagram on the right
    img_path = os.path.join(IMG_DIR, 'hinh1_er_diagram.png')
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, Inches(4.6), Inches(1.55),
                              width=Inches(8.3), height=Inches(5.0))

    # Left side: 6 tables key info
    add_text(s, 0.5, 1.6, 4.0, 0.4,
             'Tổ chức bảng',
             size=15, bold=True, color=NAVY, font='Calibri')

    # Mô tả ngắn gọn để không wrap qua nhiều dòng
    tables = [
        ('users', 'admin / staff, hash PBKDF2'),
        ('products', 'không lưu đơn giá'),
        ('warehouses', 'có is_active'),
        ('inventory_receipts', 'PENDING / APPROVED'),
        ('receipt_details', 'FK PROTECT product'),
        ('inventory', 'UNIQUE(wh, product)'),
    ]
    for i, (name, desc) in enumerate(tables):
        y = 2.1 + i * 0.55
        # Mở rộng pill thành 2.0" để chứa được "inventory_receipts" (18 char) ở 10pt Consolas
        add_round(s, 0.5, y, 2.0, 0.42, MID, corner=0.3)
        add_text(s, 0.5, y, 2.0, 0.42,
                 name, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Consolas')
        # Dồn desc qua phải, dùng font nhỏ hơn để vừa cột
        add_text(s, 2.6, y + 0.05, 1.95, 0.4,
                 desc, size=10, color=DARK_GRAY, font='Calibri')

    # Bottom callout
    add_round(s, 0.5, 6.4, 12.3, 0.7, NAVY, corner=0.05)
    add_text(s, 0.7, 6.5, 12, 0.5, [
        'Điểm chính: tách header phiếu (inventory_receipts) khỏi tồn kho thực (inventory). '
        'Tồn kho CHỈ cập nhật khi phiếu được duyệt — bảo vệ bằng transaction.atomic.'
    ], size=12, color=WHITE, italic=True, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_functions(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '6. Chức năng hệ thống',
                  '4 nhóm chính tương ứng 4 app Django')

    groups = [
        ('1', 'Xác thực & Tài khoản', 'register · login · logout · me',
         '4 endpoint', MID),
        ('2', 'Quản lý Danh mục', 'CRUD products · CRUD warehouses',
         '5 endpoint · IsAdminOrReadOnly', MID),
        ('3', 'Vận hành Kho', 'tạo phiếu · thêm chi tiết · duyệt / hủy phiếu',
         '4 endpoint · atomic transaction', ACCENT),
        ('4', 'Báo cáo', 'tồn kho theo kho · cảnh báo tồn thấp',
         '2 endpoint', MID),
    ]
    box_w = 6.0
    box_h = 2.4
    positions = [(0.5, 1.6), (6.85, 1.6), (0.5, 4.15), (6.85, 4.15)]

    for (num, head, items, badge, color), (x, y) in zip(groups, positions):
        # Outer rounded box
        add_round(s, x, y, box_w, box_h, OFFWHITE, line=LIGHT, corner=0.03)
        # Number circle
        add_circle(s, x + 0.55, y + 0.55, 0.7, color)
        add_text(s, x + 0.2, y + 0.2, 0.7, 0.7,
                 num, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Calibri')
        # Title
        add_text(s, x + 1.1, y + 0.2, box_w - 1.3, 0.5,
                 head, size=17, bold=True, color=NAVY, font='Calibri')
        # Items
        add_text(s, x + 1.1, y + 0.75, box_w - 1.3, 0.5,
                 items, size=12, color=DARK_GRAY, font='Calibri')
        # Badge bottom
        add_round(s, x + 0.3, y + box_h - 0.55, box_w - 0.6, 0.42, color, corner=0.3)
        add_text(s, x + 0.3, y + box_h - 0.55, box_w - 0.6, 0.42,
                 badge, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Calibri')

    # Note bottom
    add_text(s, 0.5, 6.7, 12.3, 0.4,
             'Tổng cộng: 14 URL endpoint, 23 thao tác HTTP, 4 app domain tách biệt.',
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_endpoints(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '7. Danh sách 14 URL endpoint REST',
                  '23 thao tác HTTP (method + URL combo)')

    rows = [
        ('POST', '/api/auth/register', 'Đăng ký', 'Public'),
        ('POST', '/api/auth/login', 'Đăng nhập, trả token', 'Public'),
        ('POST', '/api/auth/logout', 'Đăng xuất', 'Auth'),
        ('GET / PUT', '/api/auth/me', 'Xem / cập nhật profile', 'Auth'),
        ('GET / POST', '/api/products', 'List / tạo SP', 'Auth / Admin'),
        ('GET / PUT / DELETE', '/api/products/<id>', 'Xem / sửa / xoá SP', 'Auth / Admin'),
        ('GET / POST', '/api/warehouses', 'List / tạo kho', 'Auth / Admin'),
        ('GET / PUT / DELETE', '/api/warehouses/<id>', 'Xem / sửa / xoá kho', 'Auth / Admin'),
        ('GET', '/api/warehouses/<id>/inventory', 'Xem tồn kho theo kho', 'Auth'),
        ('GET / POST', '/api/receipts', 'List / tạo phiếu', 'Auth'),
        ('GET / DELETE', '/api/receipts/<id>', 'Xem / xoá phiếu PENDING', 'Auth / Admin'),
        ('POST', '/api/receipts/<id>/details', 'Thêm chi tiết phiếu', 'Auth'),
        ('PUT', '/api/receipts/<id>/status', 'Duyệt / hủy phiếu', 'Admin'),
        ('GET', '/api/reports/low-stock', 'Báo cáo tồn thấp', 'Auth'),
    ]

    # Header
    cols_x = [0.5, 3.4, 7.4, 10.7]
    cols_w = [2.85, 3.95, 3.25, 2.1]
    headers = ['HTTP Method', 'URL', 'Mô tả', 'Quyền']
    header_y = 1.55
    header_h = 0.45
    add_rect(s, 0.5, header_y, 12.3, header_h, NAVY)
    for x, w, h in zip(cols_x, cols_w, headers):
        add_text(s, x + 0.1, header_y + 0.05, w - 0.1, header_h - 0.05,
                 h, size=12, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, font='Calibri')

    # Rows
    row_y = header_y + header_h
    row_h = 0.35
    for i, row in enumerate(rows):
        bg = OFFWHITE if i % 2 == 0 else WHITE
        add_rect(s, 0.5, row_y + i * row_h, 12.3, row_h, bg)
        for j, (x, w, val) in enumerate(zip(cols_x, cols_w, row)):
            # Color method column
            color = ACCENT if j == 0 else DARK_GRAY
            weight = True if j == 1 else False
            font = 'Consolas' if j in (0, 1) else 'Calibri'
            add_text(s, x + 0.1, row_y + i * row_h + 0.04, w - 0.1, row_h - 0.05,
                     val, size=10, bold=weight, color=color,
                     anchor=MSO_ANCHOR.MIDDLE, font=font)

    add_text(s, 0.5, 7.0, 12.3, 0.3,
             'Tài liệu hóa tự động qua Swagger UI tại /swagger/  (drf-yasg)',
             size=11, italic=True, color=MID, align=PP_ALIGN.CENTER, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_demo_atomic(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '8. Demo nghiệp vụ then chốt: Duyệt phiếu xuất',
                  'Bảo vệ tồn kho bằng transaction.atomic + select_for_update')

    # Bên trái: code snippet
    add_text(s, 0.5, 1.6, 6.5, 0.35,
             'Đoạn code đặc trưng (inventory/views.py)',
             size=13, bold=True, color=NAVY, font='Calibri')

    code_bg = add_round(s, 0.5, 1.95, 6.5, 4.65, hex_rgb('#1E1E1E'), corner=0.02)
    code = (
        'with transaction.atomic():\n'
        '    for ct in chi_tiet_list:\n'
        '        ton, _ = (Inventory.objects\n'
        '          .select_for_update()\n'
        '          .get_or_create(\n'
        '            warehouse_id=phieu.warehouse_id,\n'
        '            product_id=ct.product_id,\n'
        '            defaults={"quantity": 0},\n'
        '          ))\n'
        '        if phieu.receipt_type == IMPORT:\n'
        '            ton.quantity += ct.quantity\n'
        '        else:\n'
        '            ton.quantity -= ct.quantity\n'
        '        ton.save()\n'
        '    phieu.receipt_status = APPROVED\n'
        '    phieu.save()'
    )
    add_text(s, 0.7, 2.1, 6.2, 4.4,
             code.split('\n'),
             size=11, color=hex_rgb('#E0E0E0'), font='Consolas')

    # Bên phải: 3 bước nghiệp vụ
    add_text(s, 7.4, 1.6, 5.4, 0.35,
             '3 bước bảo vệ dữ liệu',
             size=13, bold=True, color=NAVY, font='Calibri')

    steps = [
        ('1', 'Kiểm tra tồn ĐỦ',
         'Nếu xuất quá tồn hiện có → trả 400 ngay, KHÔNG đụng DB.'),
        ('2', 'Khóa hàng (row-level lock)',
         'select_for_update() bắt đầu transaction, khóa dòng tồn kho cho '
         'tới khi commit. Admin khác duyệt cùng SP phải chờ.'),
        ('3', 'Cộng/trừ tồn + đổi trạng thái',
         'Cập nhật atomic — nếu fail giữa chừng, rollback toàn bộ.'),
    ]
    for i, (num, head, body) in enumerate(steps):
        y = 2.1 + i * 1.55
        add_round(s, 7.4, y, 5.4, 1.4, OFFWHITE, line=LIGHT, corner=0.05)
        add_circle(s, 7.8, y + 0.35, 0.6, ACCENT)
        add_text(s, 7.45, y + 0.05, 0.7, 0.6,
                 num, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Calibri')
        add_text(s, 8.3, y + 0.05, 4.5, 0.35,
                 head, size=13, bold=True, color=NAVY, font='Calibri')
        add_text(s, 8.3, y + 0.4, 4.5, 0.9,
                 body, size=11, color=DARK_GRAY, font='Calibri')

    # Bottom highlight
    add_round(s, 0.5, 6.7, 12.3, 0.4, NAVY, corner=0.1)
    add_text(s, 0.5, 6.72, 12.3, 0.4,
             'Race condition: hai admin duyệt cùng SP cùng lúc → tuần tự hoá '
             '— KHÔNG để tồn kho âm.',
             size=11, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_testing(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '9. Kiểm thử',
                  'Postman collection + Swagger UI live')

    # Bên trái: Postman card
    add_round(s, 0.5, 1.6, 6.0, 5.3, OFFWHITE, line=LIGHT, corner=0.03)
    add_round(s, 0.5, 1.6, 6.0, 0.55, ACCENT, corner=0.1)
    add_text(s, 0.7, 1.65, 5.6, 0.45,
             'POSTMAN COLLECTION', size=14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, font='Calibri')

    add_text(s, 0.7, 2.3, 5.6, 0.4,
             'ezware_postman_collection.json',
             size=12, italic=True, color=NAVY, font='Consolas')

    bullets = [
        '30+ test case bao phủ 23 thao tác',
        'Script JS tự lưu token sau khi login',
        'Token được auto-inject cho mọi request sau',
        'Có test happy path & error path',
        'Regression test cho lỗ hổng leo thang quyền (BUG #1)',
    ]
    for i, b in enumerate(bullets):
        y = 2.8 + i * 0.55
        add_circle(s, 0.95, y + 0.1, 0.18, MID)
        add_text(s, 1.2, y, 5.2, 0.5,
                 b, size=12, color=DARK_GRAY, font='Calibri')

    # Bên phải: Swagger UI
    add_round(s, 6.85, 1.6, 6.0, 5.3, OFFWHITE, line=LIGHT, corner=0.03)
    add_round(s, 6.85, 1.6, 6.0, 0.55, NAVY, corner=0.1)
    add_text(s, 7.05, 1.65, 5.6, 0.45,
             'SWAGGER UI (drf-yasg)', size=14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, font='Calibri')

    # Insert Swagger UI screenshot — chỉ set width, để python-pptx auto-fit aspect ratio
    # (ảnh gốc 14:17 narrow tall — set width 4.0" → height auto ~4.85")
    img_path = os.path.join(IMG_DIR, 'hinh6_swagger_ui.png')
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, Inches(7.85), Inches(2.25),
                              width=Inches(4.0))

    add_slide_number_footer(s, slide_num, total)


def build_deploy(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '10. Deploy thực tế lên Render',
                  'Hệ thống đã chạy công khai trên Internet')

    # URL banner top
    add_round(s, 0.5, 1.6, 12.3, 0.85, NAVY, corner=0.05)
    add_text(s, 0.7, 1.65, 12, 0.4,
             'URL công khai (người chấm có thể truy cập trực tiếp):',
             size=11, italic=True, color=LIGHT, font='Calibri')
    add_text(s, 0.7, 1.95, 12, 0.5,
             f'{DEPLOY_URL}/swagger/',
             size=20, bold=True, color=ACCENT, font='Consolas')

    # 3 bước deploy
    steps = [
        ('1', 'PUSH', 'Push source lên GitHub',
         'requirements.txt (5 package production)'),
        ('2', 'BUILD', 'Render build tự động',
         'pip install · collectstatic · migrate'),
        ('3', 'SERVE', 'Gunicorn + WhiteNoise serve',
         '23 endpoint REST + Swagger UI'),
    ]
    sw = 4.0
    for i, (num, head, sub, body) in enumerate(steps):
        x = 0.5 + i * (sw + 0.15)
        y = 2.7
        add_round(s, x, y, sw, 1.7, OFFWHITE, line=LIGHT, corner=0.05)
        add_circle(s, x + 0.5, y + 0.45, 0.6, ACCENT)
        add_text(s, x + 0.2, y + 0.15, 0.6, 0.6,
                 num, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Calibri')
        add_text(s, x + 1.2, y + 0.1, sw - 1.3, 0.35,
                 head, size=15, bold=True, color=NAVY, font='Calibri')
        add_text(s, x + 1.2, y + 0.45, sw - 1.3, 0.35,
                 sub, size=11, italic=True, color=MID, font='Calibri')
        add_text(s, x + 0.25, y + 0.95, sw - 0.4, 0.65,
                 body, size=11, color=DARK_GRAY, font='Calibri')

    # Bottom: env vars + technical highlights
    add_text(s, 0.5, 4.55, 12.3, 0.35,
             'Cấu hình production',
             size=14, bold=True, color=NAVY, font='Calibri')

    items = [
        ('SECRET_KEY', 'random 50 ký tự từ env var (không hardcode)'),
        ('DEBUG', 'False ở production, True ở local (đọc từ env)'),
        ('ALLOWED_HOSTS', f'ezware-backend.onrender.com'),
        ('Static files', 'WhiteNoise nén + cache-bust Swagger UI assets'),
        ('Database', 'SQLite (đủ cho demo) — có thể swap Postgres qua DATABASE_URL'),
    ]
    for i, (k, v) in enumerate(items):
        y = 5.0 + i * 0.38
        add_round(s, 0.5, y, 4.0, 0.32, NAVY, corner=0.3)
        add_text(s, 0.5, y, 4.0, 0.32,
                 k, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Consolas')
        add_text(s, 4.7, y, 8.1, 0.32,
                 v, size=11, color=DARK_GRAY,
                 anchor=MSO_ANCHOR.MIDDLE, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_challenges(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, '11. Thách thức kỹ thuật đã giải quyết',
                  'Học hỏi qua bug thực tế')

    challenges = [
        ('OOM trên Render free tier',
         'pip install pandas + numpy + Django đồng thời > 512 MB RAM → SIGKILL silent.',
         'Tách requirements.txt: bỏ pandas/numpy ra requirements-dev.txt. Build giảm 50% memory.'),
        ('pkg_resources biến mất',
         'setuptools 81+ loại bỏ pkg_resources khỏi wheel. drf-yasg 1.21.7 vẫn import → ModuleNotFoundError.',
         'Pin setuptools<81 → pip pull 80.10.2 (phiên bản cuối còn pkg_resources).'),
        ('Privilege escalation qua register',
         'API /auth/register cho phép client gửi user_role="admin" → ai cũng tự thăng cấp admin.',
         'Bỏ user_role ra khỏi serializer; force STAFF khi register; viết regression test.'),
        ('Race condition khi duyệt phiếu',
         '2 admin duyệt cùng SP cùng lúc có thể làm tồn kho âm.',
         'Bọc trong transaction.atomic + select_for_update để khóa hàng tồn kho.'),
    ]
    for i, (title, problem, solution) in enumerate(challenges):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.35
        y = 1.55 + row * 2.65

        add_round(s, x, y, 6.0, 2.45, OFFWHITE, line=LIGHT, corner=0.03)
        # Header
        add_round(s, x, y, 6.0, 0.45, ACCENT, corner=0.1)
        add_text(s, x + 0.2, y + 0.05, 5.6, 0.35,
                 title, size=13, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, font='Calibri')
        # Problem
        add_text(s, x + 0.2, y + 0.55, 5.6, 0.3,
                 'VẤN ĐỀ:', size=10, bold=True, color=NAVY, font='Calibri')
        add_text(s, x + 0.2, y + 0.85, 5.6, 0.75,
                 problem, size=11, color=DARK_GRAY, font='Calibri')
        # Solution
        add_text(s, x + 0.2, y + 1.55, 5.6, 0.3,
                 'GIẢI PHÁP:', size=10, bold=True, color=MID, font='Calibri')
        add_text(s, x + 0.2, y + 1.85, 5.6, 0.55,
                 solution, size=11, color=DARK_GRAY, font='Calibri')

    add_slide_number_footer(s, slide_num, total)


def build_conclusion(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, NAVY)

    # Decorative circle
    add_circle(s, 11.5, 0.8, 1.5, MID)
    add_circle(s, 12.5, 0.3, 0.6, ACCENT)

    add_text(s, 0.5, 0.5, 12.3, 0.7,
             '12. Kết luận',
             size=36, bold=True, color=WHITE, font='Calibri')
    add_text(s, 0.5, 1.2, 12.3, 0.4,
             'Đã đạt — Hài lòng nhất — Hạn chế — Hướng phát triển',
             size=15, italic=True, color=LIGHT, font='Calibri')

    # 4 cột
    items = [
        ('ĐÃ ĐẠT', ACCENT, [
            'Backend REST API hoàn chỉnh',
            '14 URL · 23 ops · 6 bảng',
            'Phân quyền 2 vai trò',
            'Deploy công khai',
            '30+ Postman test case',
        ]),
        ('HÀI LÒNG NHẤT', MID, [
            'Nghiệp vụ duyệt phiếu xuất',
            'với transaction.atomic',
            '+ select_for_update',
            'chống được race condition',
            'khi nhiều admin duyệt cùng lúc',
        ]),
        ('HẠN CHẾ', hex_rgb('#A0A0A0'), [
            'SQLite — chưa scale lớn',
            'Chưa có frontend riêng',
            'Render free → service ngủ sau 15 phút',
            'Chưa hỗ trợ chuyển kho trực tiếp',
            'Chưa có MFA / refresh token',
        ]),
        ('HƯỚNG PHÁT TRIỂN', LIGHT, [
            'Postgres + dj-database-url',
            'Frontend React/Vue',
            'JWT thay Token đơn',
            'Logging + Sentry',
            'Báo cáo doanh thu / nhập-xuất',
        ]),
    ]

    col_w = 3.0
    for i, (head, color, lines) in enumerate(items):
        x = 0.5 + i * (col_w + 0.1)
        y = 2.0
        add_round(s, x, y, col_w, 4.7, hex_rgb('#FFFFFF'), corner=0.05)
        add_round(s, x, y, col_w, 0.55, color, corner=0.1)
        # head color check — set text color for readability
        head_fg = WHITE if color in (ACCENT, MID, NAVY) else NAVY
        add_text(s, x, y + 0.07, col_w, 0.4,
                 head, size=13, bold=True, color=head_fg,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Calibri')
        bullet_text = ['•  ' + l for l in lines]
        add_text(s, x + 0.2, y + 0.7, col_w - 0.3, 3.9,
                 bullet_text, size=11, color=DARK_GRAY, font='Calibri')

    # Bottom URL
    add_text(s, 0.5, 7.0, 12.3, 0.35,
             f'Demo live: {DEPLOY_URL}/swagger/',
             size=14, italic=True, color=ACCENT,
             align=PP_ALIGN.CENTER, font='Consolas')


def build_thanks(slide_num, total, prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, NAVY)

    # Big decorative
    add_circle(s, 12.0, 1.5, 2.5, MID)
    add_circle(s, 11.0, 0.8, 1.0, ACCENT)
    add_circle(s, 1.5, 6.5, 1.8, MID)

    add_text(s, 0.5, 2.5, 12.3, 1.1,
             'CẢM ƠN THẦY/CÔ ĐÃ LẮNG NGHE',
             size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font='Calibri')

    add_text(s, 0.5, 3.5, 12.3, 0.55,
             'Sẵn sàng nhận câu hỏi',
             size=22, italic=True, color=LIGHT,
             align=PP_ALIGN.CENTER, font='Calibri')

    add_round(s, 3.0, 4.7, 7.3, 1.7, MID, corner=0.05)
    add_text(s, 3.0, 4.8, 7.3, 0.4,
             'NHÓM XX',
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font='Calibri')
    add_text(s, 3.0, 5.2, 7.3, 1.2, [
        'Nguyễn Văn A  —  25XXXXX1',
        'Nguyễn Văn B  —  25XXXXX2',
    ], size=15, color=WHITE, align=PP_ALIGN.CENTER, font='Calibri')

    add_text(s, 0.5, 6.8, 12.3, 0.4,
             f'{DEPLOY_URL}/swagger/',
             size=16, italic=True, color=ACCENT, font='Consolas',
             align=PP_ALIGN.CENTER)


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add blank layout reference (the 7th built-in is Blank in default template)
    # Default order: 0 Title slide, 1 Title+Content, 2 Section header, 3 Two content,
    # 4 Comparison, 5 Title only, 6 Blank
    builders = [
        build_cover,
        build_agenda,
        build_intro,
        build_goals,
        build_tech,
        build_architecture,
        build_database,
        build_functions,
        build_endpoints,
        build_demo_atomic,
        build_testing,
        build_deploy,
        build_challenges,
        build_conclusion,
        build_thanks,
    ]
    total = len(builders)
    for i, b in enumerate(builders, start=1):
        b(i, total, prs)

    prs.save(OUTPUT)
    print(f'Đã sinh {total} slide: {OUTPUT}')
    print(f'File size: {os.path.getsize(OUTPUT) / 1024:.1f} KB')


if __name__ == '__main__':
    main()
